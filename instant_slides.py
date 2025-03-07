from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import json
import os

# 🎨 Color Palette
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)
ACCENT_BLUE = RGBColor(30, 136, 229)  # Deep Blue
ACCENT_ORANGE = RGBColor(255, 87, 34)  # Orange

# Default slide layout mapping
LAYOUTS = {
    "title": 0,  # Title Slide
    "text": 1,   # Title + Content
    "image": 5   # Title + Image
}

def create_presentation(slides_data, output_file="generated_slides.pptx"):
    """Generate a Presentation Zen-style PowerPoint."""
    prs = Presentation()

    for idx, slide in enumerate(slides_data):
        slide_type = slide.get("type", "text")
        layout_index = LAYOUTS.get(slide_type, 1)  # Default to text layout

        ppt_slide = prs.slides.add_slide(prs.slide_layouts[layout_index])
        title_shape = ppt_slide.shapes.title
        content_box = ppt_slide.placeholders[1] if layout_index == 1 else None

        # 🌟 Apply Clean White Background
        ppt_slide.background.fill.solid()
        ppt_slide.background.fill.fore_color.rgb = WHITE

        # 📝 Set Title
        if title_shape and "title" in slide:
            title_shape.text = slide["title"]
            title_shape.text_frame.paragraphs[0].font.size = Pt(44)  # Large font
            title_shape.text_frame.paragraphs[0].font.bold = True
            title_shape.text_frame.paragraphs[0].font.color.rgb = BLACK

        # 📌 Set Minimal Text Content
        if content_box and "text" in slide:
            content_box.text = slide["text"]
            content_box.text_frame.paragraphs[0].font.size = Pt(32)  # Big, readable text
            content_box.text_frame.paragraphs[0].font.color.rgb = BLACK

        # 🖼 Insert Image (Full Width)
        if "image" in slide and layout_index == 5:
            img_path = slide["image"]
            if os.path.exists(img_path):
                ppt_slide.shapes.add_picture(img_path, Inches(0.5), Inches(1), width=Inches(9))
            else:
                print(f"⚠️ Image not found: {img_path}")

        # 🔥 Add Optional Accent Color to Titles
        if idx % 2 == 0:  # Alternate colors for visual variety
            title_shape.text_frame.paragraphs[0].font.color.rgb = ACCENT_BLUE
        else:
            title_shape.text_frame.paragraphs[0].font.color.rgb = ACCENT_ORANGE

    prs.save(output_file)
    print(f"✅ Presentation saved as {output_file}")

# Example JSON input
example_slides = [
    {"type": "title", "title": "Welcome to Orchestrate"},
    {"type": "text", "title": "Less is More", "text": "One big idea per slide. No clutter."},
    {"type": "image", "title": "Visual Storytelling", "image": "example_image.jpg"}
]

# Generate PPTX
if __name__ == "__main__":
    create_presentation(example_slides)
